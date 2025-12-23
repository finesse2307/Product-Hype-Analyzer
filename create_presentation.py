#!/usr/bin/env python3
"""
Create PowerPoint presentation for viral products analysis
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Check if visualizations exist
if not os.path.exists('visualizations/00_dashboard.png'):
    print("ERROR: Visualizations not found. Run create_visualizations.py first.")
    exit(1)

print("=" * 80)
print("CREATING POWERPOINT PRESENTATION")
print("=" * 80)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define colors
DARK_BLUE = RGBColor(46, 134, 171)      # #2E86AB
PURPLE = RGBColor(162, 59, 114)         # #A23B72  
ORANGE = RGBColor(241, 143, 1)          # #F18F01
RED = RGBColor(199, 62, 29)             # #C73E1D

# ============================================================================
# Slide 1: Title Slide
# ============================================================================
print("\n1. Creating title slide...")

blank_layout = prs.slide_layouts[6]  # Blank layout
slide = prs.slides.add_slide(blank_layout)

# Add dark blue background
background = slide.shapes.add_shape(
    1,  # Rectangle
    0, 0, prs.slide_width, prs.slide_height
)
background.fill.solid()
background.fill.fore_color.rgb = DARK_BLUE
background.line.fill.background()

# Main title
title_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
)
title_frame = title_box.text_frame
title_frame.text = "Viral Product Hype Cycle"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(54)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(255, 255, 255)
title_para.alignment = PP_ALIGN.CENTER

# Subtitle
subtitle_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(4), Inches(9), Inches(0.8)
)
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "SQL Analysis of Social Media Trends & Consumer Products"
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.font.size = Pt(24)
subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
subtitle_para.alignment = PP_ALIGN.CENTER

# Footer
footer_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(6.8), Inches(9), Inches(0.5)
)
footer_frame = footer_box.text_frame
footer_frame.text = "Stanley • Owala • Hydro Flask • YETI | 2020-2024 Analysis"
footer_para = footer_frame.paragraphs[0]
footer_para.font.size = Pt(14)
footer_para.font.italic = True
footer_para.font.color.rgb = RGBColor(255, 255, 255)
footer_para.alignment = PP_ALIGN.CENTER

# ============================================================================
# Slide 2: Dashboard Overview
# ============================================================================
print("2. Creating dashboard overview...")

slide = prs.slides.add_slide(blank_layout)

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "Executive Dashboard"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(32)
title_para.font.bold = True
title_para.font.color.rgb = DARK_BLUE

# Add dashboard image
slide.shapes.add_picture(
    'visualizations/00_dashboard.png',
    Inches(0.3), Inches(1), 
    width=Inches(9.4)
)

# ============================================================================
# Slide 3: Search Interest Trends
# ============================================================================
print("3. Creating trends slide...")

slide = prs.slides.add_slide(blank_layout)

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "Search Interest Over Time"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(32)
title_para.font.bold = True
title_para.font.color.rgb = DARK_BLUE

# Add trend chart
slide.shapes.add_picture(
    'visualizations/01_trend_lines.png',
    Inches(0.5), Inches(1.1),
    width=Inches(9)
)

# Key insight box
insight_box = slide.shapes.add_textbox(Inches(6.5), Inches(6.3), Inches(3), Inches(0.9))
insight_box.fill.solid()
insight_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
insight_box.line.color.rgb = ORANGE
insight_box.line.width = Pt(2)

insight_frame = insight_box.text_frame
insight_frame.text = "Stanley peaked March 2023\nOwala peaked July 2023"
for para in insight_frame.paragraphs:
    para.font.size = Pt(14)
    para.font.bold = True
    para.font.color.rgb = DARK_BLUE

# ============================================================================
# Slide 4: Peak Comparison
# ============================================================================
print("4. Creating peak comparison...")

slide = prs.slides.add_slide(blank_layout)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "Peak vs Average Interest"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(32)
title_para.font.bold = True
title_para.font.color.rgb = DARK_BLUE

slide.shapes.add_picture(
    'visualizations/02_peak_comparison.png',
    Inches(1.5), Inches(1.2),
    width=Inches(7)
)

# Insight
insight_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(8.5), Inches(0.9))
insight_box.fill.solid()
insight_box.fill.fore_color.rgb = RGBColor(250, 250, 250)

insight_frame = insight_box.text_frame
insight_frame.text = "Products maintain 15-20% of peak interest after viral cycle completes"
insight_para = insight_frame.paragraphs[0]
insight_para.font.size = Pt(16)
insight_para.font.color.rgb = DARK_BLUE

# ============================================================================
# Slide 5: Market Share Evolution
# ============================================================================
print("5. Creating market share slide...")

slide = prs.slides.add_slide(blank_layout)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "Market Share Competition"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(32)
title_para.font.bold = True
title_para.font.color.rgb = DARK_BLUE

slide.shapes.add_picture(
    'visualizations/03_market_share.png',
    Inches(0.5), Inches(1.1),
    width=Inches(9)
)

insight_box = slide.shapes.add_textbox(Inches(6.2), Inches(6.3), Inches(3.3), Inches(0.9))
insight_box.fill.solid()
insight_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
insight_box.line.color.rgb = PURPLE
insight_box.line.width = Pt(2)

insight_frame = insight_box.text_frame
insight_frame.text = "Owala captured 25.7% share\nin July 2024, leading market"
for para in insight_frame.paragraphs:
    para.font.size = Pt(14)
    para.font.bold = True
    para.font.color.rgb = DARK_BLUE

# ============================================================================
# Slide 6: Event Impact
# ============================================================================
print("6. Creating event impact slide...")

slide = prs.slides.add_slide(blank_layout)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "Retail Event Impact Analysis"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(32)
title_para.font.bold = True
title_para.font.color.rgb = DARK_BLUE

slide.shapes.add_picture(
    'visualizations/04_event_impact.png',
    Inches(0.8), Inches(1.1),
    width=Inches(8.5)
)

# ============================================================================
# Slide 7: TikTok Engagement
# ============================================================================
print("7. Creating TikTok engagement slide...")

slide = prs.slides.add_slide(blank_layout)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "TikTok Engagement & Sentiment"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(32)
title_para.font.bold = True
title_para.font.color.rgb = DARK_BLUE

slide.shapes.add_picture(
    'visualizations/05_tiktok_trends.png',
    Inches(0.5), Inches(1.1),
    width=Inches(9)
)

# ============================================================================
# Slide 8: Growth Heatmap
# ============================================================================
print("8. Creating growth heatmap...")

slide = prs.slides.add_slide(blank_layout)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "Growth Momentum Patterns"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(32)
title_para.font.bold = True
title_para.font.color.rgb = DARK_BLUE

slide.shapes.add_picture(
    'visualizations/06_growth_heatmap.png',
    Inches(0.5), Inches(1.3),
    width=Inches(9)
)

insight_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(8.5), Inches(0.9))
insight_box.fill.solid()
insight_box.fill.fore_color.rgb = RGBColor(250, 250, 250)

insight_frame = insight_box.text_frame
insight_frame.text = "Green = Growth periods | Red = Decline | Peak growth: 800% week-over-week"
insight_para = insight_frame.paragraphs[0]
insight_para.font.size = Pt(16)
insight_para.font.color.rgb = DARK_BLUE

# ============================================================================
# Slide 9: Key Findings
# ============================================================================
print("9. Creating key findings slide...")

slide = prs.slides.add_slide(blank_layout)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "Key Findings"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = DARK_BLUE

# Findings content
findings_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(5.5))
findings_frame = findings_box.text_frame
findings_frame.word_wrap = True

findings = [
    "Viral products peak 3-6 months after initial social media surge",
    "Maximum growth rates of 800% week-over-week during viral phase",
    "Target sellouts drive average +15% search interest lift",
    "Products stabilize at 15-20% of peak interest after maturity",
    "Market leadership shifts: Owala overtook Stanley in mid-2024",
    "TikTok engagement rates average 2.5-3.5% (high for consumer products)",
    "Sentiment remains positive (0.58-0.72) throughout lifecycle",
]

for i, finding in enumerate(findings):
    p = findings_frame.add_paragraph()
    p.text = finding
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.space_before = Pt(8)
    p.level = 0
    p.bullet = True

# ============================================================================
# Slide 10: SQL Techniques
# ============================================================================
print("10. Creating SQL techniques slide...")

slide = prs.slides.add_slide(blank_layout)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "SQL Techniques Demonstrated"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = DARK_BLUE

# Left column
left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(4.2), Inches(5.5))
left_frame = left_box.text_frame
left_frame.word_wrap = True

left_techniques = [
    "Window Functions (LAG, LEAD, ROW_NUMBER)",
    "Common Table Expressions (CTEs)",
    "Complex Multi-Table Joins",
    "Date Calculations (JULIANDAY, strftime)",
    "Event-Based Windowing",
]

for technique in left_techniques:
    p = left_frame.add_paragraph()
    p.text = technique
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.space_before = Pt(8)
    p.bullet = True

# Right column
right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.2), Inches(5.5))
right_frame = right_box.text_frame
right_frame.word_wrap = True

right_techniques = [
    "Percentage & Growth Calculations",
    "Market Share Analysis",
    "Rolling Averages",
    "Correlation Analysis",
    "Statistical Aggregations",
]

for technique in right_techniques:
    p = right_frame.add_paragraph()
    p.text = technique
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.space_before = Pt(8)
    p.bullet = True

# ============================================================================
# Slide 11: Conclusion
# ============================================================================
print("11. Creating conclusion slide...")

slide = prs.slides.add_slide(blank_layout)

# Background
background = slide.shapes.add_shape(
    1, 0, 0, prs.slide_width, prs.slide_height
)
background.fill.solid()
background.fill.fore_color.rgb = RGBColor(245, 245, 245)
background.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "Business Applications"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(40)
title_para.font.bold = True
title_para.font.color.rgb = DARK_BLUE
title_para.alignment = PP_ALIGN.CENTER

content_box = slide.shapes.add_textbox(Inches(1.5), Inches(3), Inches(7), Inches(3))
content_frame = content_box.text_frame
content_frame.word_wrap = True

applications = [
    "Inventory Planning: Forecast demand spikes 2-3 weeks ahead",
    "Marketing Timing: Schedule limited editions during growth phase",
    "Investment Strategy: Identify entry points before viral peaks",
    "Competitive Intelligence: Track market share shifts in real-time",
]

for app in applications:
    p = content_frame.add_paragraph()
    p.text = app
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.space_before = Pt(12)
    p.bullet = True

# ============================================================================
# Save presentation
# ============================================================================
output_file = 'Viral_Products_Analysis.pptx'
prs.save(output_file)

print(f"\n✓ Presentation saved: {output_file}")

print("\n" + "=" * 80)
print("✓ PRESENTATION COMPLETE!")
print("=" * 80)
print(f"\nCreated: {output_file}")
print("Slides: 11")
print("  1. Title Slide")
print("  2. Executive Dashboard")
print("  3. Search Interest Trends")
print("  4. Peak Comparison")
print("  5. Market Share Evolution")
print("  6. Retail Event Impact")
print("  7. TikTok Engagement")
print("  8. Growth Heatmap")
print("  9. Key Findings")
print("  10. SQL Techniques")
print("  11. Business Applications")
